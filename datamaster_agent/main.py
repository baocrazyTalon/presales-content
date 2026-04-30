"""
main.py
────────
Talon.One Presales Multi-Agent System — Entry Point

Usage:
    # Interactive CLI
    python main.py

    # Headless (pass a JSON prospect file)
    python main.py --prospect prospect.json

    # Start the FastAPI server
    python main.py --serve

Environment:
    Copy .env.example → .env and fill in your keys before running.
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import sys

from dotenv import load_dotenv

load_dotenv()

# ── Logging setup ─────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(name)s — %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("presales_agent")

# ── Lazy import — avoids crashes if deps aren't installed yet ─────────────────
def _import_graph():
    from core.graph import get_compiled_graph
    return get_compiled_graph()

from core.state import AgentState, ProspectContext
from core.memory import make_thread_id, get_thread_history


# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────

def build_initial_state(prospect: ProspectContext, thread_id: str = "") -> AgentState:
    """Construct a clean AgentState for a new presales run."""
    return AgentState(
        messages=[],
        prospect=prospect,
        tasks=[],
        completed_tasks=[],
        rag_queries=[],
        retrieved_chunks=[],
        selling_agenda=[],
        narrative_arc=None,
        key_talking_points=[],
        competitor_kill_shots=[],
        html_output=None,
        diagram_definitions=[],
        validation=None,
        output_filename=None,
        github_pr_url=None,
        vercel_url=None,
        current_agent="coordinator",
        iteration_count=0,
        error=None,
        thread_id=thread_id,
        personal_insights=[],
        user_feedback=None,
    )


def run_agent(prospect: ProspectContext, resume: bool = False) -> AgentState:
    """Invoke the compiled LangGraph and return the final state."""
    graph = _import_graph()
    thread_id = make_thread_id(prospect.get("company_name", "default"))
    initial = build_initial_state(prospect, thread_id=thread_id)
    config = {"configurable": {"thread_id": thread_id}}

    logger.info(
        "Starting presales agent for %s (%s) thread=%s resume=%s",
        prospect.get("company_name", "Unknown"),
        prospect.get("industry", "Unknown"),
        thread_id,
        resume,
    )

    if resume:
        # Resume from last checkpoint — graph.invoke with an existing
        # thread_id will pick up from the last saved state.
        final_state: AgentState = graph.invoke(None, config=config)
    else:
        final_state = graph.invoke(initial, config=config)

    vercel = final_state.get("vercel_url")
    if vercel:
        logger.info("✅  Deployment complete → %s", vercel)
    elif final_state.get("error"):
        logger.error("❌  Agent stopped with error: %s", final_state["error"])
    else:
        logger.warning("⚠️   Agent finished without deploying")

    return final_state


# ─────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────

def _prompt_prospect() -> ProspectContext:
    """Interactive prompt to collect prospect details."""
    print("\n── Talon.One Presales Agent ─────────────────────────────")
    print("Answer the questions below, then press Enter to generate.\n")

    company   = input("Company name:               ").strip()
    industry  = input("Industry (e.g. Retail/QSR): ").strip()
    use_case  = input("Use-case / goal:            ").strip()
    competitor = input("Competitor to displace:     ").strip()
    integrations_raw = input("Key integrations (comma-separated, e.g. Braze,Segment): ").strip()
    urls_raw  = input("URLs to scrape (comma-separated, or blank): ").strip()
    notes     = input("Paste raw notes (or leave blank): ").strip()

    return ProspectContext(
        company_name=company,
        industry=industry,
        use_case=use_case,
        competitor=competitor or None,
        integrations=[i.strip() for i in integrations_raw.split(",") if i.strip()],
        urls_to_scrape=[u.strip() for u in urls_raw.split(",") if u.strip()],
        raw_notes=notes,
    )


# ─────────────────────────────────────────────────────────────
# FastAPI server — app exposed at module level for uvicorn
# ─────────────────────────────────────────────────────────────

import pathlib
import re
from datetime import datetime

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

app = FastAPI(title="Talon.One Presales Agent", version="0.1.0")

# Allow requests from file:// and any localhost origin (browser wizard)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve the whole Documents folder as static files at /docs
_DOCS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
_DOCS_ROOT = pathlib.Path(_DOCS_DIR)
app.mount("/docs", StaticFiles(directory=_DOCS_DIR, html=True), name="docs")


# ── Request / Response models (module-level so FastAPI can introspect them) ──

class ProspectRequest(BaseModel):
    company_name: str
    industry: str
    use_case: str
    competitor: str | None = None
    integrations: list[str] = []
    raw_notes: str = ""
    urls_to_scrape: list[str] = []
    presentation_type: str = "client"       # client | competition | partner
    presentation_category: str = ""         # e.g. LOYALTY, CEP, CDP
    template_type: str = "sales-pitch"       # sales-pitch | technical-integration | business-case | battle-card | poc-scope | partner-gtm
    presentation_name: str = ""              # custom title / filename; auto-generated by backend if blank
    brand: dict | None = None                # {name, country, website, primary_color, secondary_color, logo_url, logo_b64}

class RunResponse(BaseModel):
    vercel_url: str | None = None
    github_pr_url: str | None = None
    output_filename: str | None = None
    local_path: str | None = None           # relative to Documents root, e.g. CLIENTS/RUSTAN/rustan-presales.html
    validation_passed: bool = False
    error: str | None = None

class ChatRequest(BaseModel):
    doc_path: str                           # relative path from Documents root
    message: str
    history: list[dict] = []

class ChatResponse(BaseModel):
    reply: str
    updated_html: str | None = None         # set when the agent regenerated the doc

class ScrapeRequest(BaseModel):
    url: str
    source: str = "web"
    title: str = ""

class BrandAnalyseRequest(BaseModel):
    url: str
    brand_name: str = ""

class CrawlRequest(BaseModel):
    start_url: str
    depth: int = 2
    url_filter: str = ""
    source: str = "web"
    max_pages: int = 10

class ScrapeResponse(BaseModel):
    url: str
    chunks_ingested: int


# ── Routes ────────────────────────────────────────────────────

@app.get("/", include_in_schema=False)
def root():
    return RedirectResponse(url="/docs/index.html")

@app.get("/health")
def health():
    return {"status": "ok"}

def _slugify(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")


DRAFT_BADGE = """
<style>
#draft-badge {
  position: fixed; top: 12px; right: 12px; z-index: 9999;
  background: #ff9800; color: #fff; font-family: sans-serif;
  font-size: 12px; font-weight: 700; padding: 4px 12px;
  border-radius: 4px; letter-spacing: .05em; pointer-events: none;
}
</style>
<div id="draft-badge">DRAFT</div>
"""


def _inject_draft(html: str) -> str:
    """Inject a DRAFT badge into the HTML body."""
    if "<body" in html:
        return html.replace("<body", DRAFT_BADGE + "<body", 1)
    return DRAFT_BADGE + html


def _resolve_output_folder(ptype: str, category: str, company: str) -> pathlib.Path:
    ptype = ptype.lower()
    cat = (category or "OTHER").upper().replace(" ", "_")
    if ptype == "competition":
        return _DOCS_ROOT / "COMPETITION" / cat
    elif ptype == "partner":
        return _DOCS_ROOT / "PARTNERS" / cat
    else:  # client (default)
        return _DOCS_ROOT / "CLIENTS" / company.upper().replace(" ", "_")


@app.post("/run", response_model=RunResponse)
def run_endpoint(req: ProspectRequest):
    prospect = ProspectContext(**req.model_dump())
    state = run_agent(prospect)

    html = state.get("html_output") or ""
    local_rel: str | None = None

    if html:
        html_draft = _inject_draft(html)
        folder = _resolve_output_folder(
            req.presentation_type,
            req.presentation_category,
            req.company_name,
        )
        folder.mkdir(parents=True, exist_ok=True)
        slug = _slugify(req.company_name)
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        filepath = folder / f"{slug}-presales-{ts}.html"
        filepath.write_text(html_draft, encoding="utf-8")
        local_rel = str(filepath.relative_to(_DOCS_ROOT))
        logger.info("[run] Saved HTML → %s", filepath)

    return RunResponse(
        vercel_url=state.get("vercel_url"),
        github_pr_url=state.get("github_pr_url"),
        output_filename=state.get("output_filename"),
        local_path=local_rel,
        validation_passed=bool(state.get("validation") and state["validation"].get("passed")),
        error=state.get("error"),
    )


@app.post("/chat", response_model=ChatResponse)
def chat_endpoint(req: ChatRequest):
    """Refine a generated document via conversation with the coordinator."""
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_core.messages import HumanMessage, SystemMessage

    # Read the current document from disk
    doc_path = _DOCS_ROOT / req.doc_path
    if not doc_path.exists() or not doc_path.suffix == ".html":
        raise HTTPException(status_code=404, detail="Document not found")
    current_html = doc_path.read_text(encoding="utf-8")

    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-pro",
        google_api_key=os.getenv("GOOGLE_API_KEY", ""),
    )

    system = SystemMessage(
        content=(
            "You are a Talon.One presales expert helping refine a sales document. "
            "When the user asks for a change, either explain the improvement or, "
            "if they ask you to update the document, respond with the full updated HTML "
            "wrapped in a markdown fenced code block labelled \"html\". "
            "Otherwise respond naturally in plain text."
        )
    )

    history_msgs = []
    for h in req.history:
        role = h.get("role", "user")
        content = h.get("content", "")
        if role == "user":
            history_msgs.append(HumanMessage(content=content))
        else:
            from langchain_core.messages import AIMessage
            history_msgs.append(AIMessage(content=content))

    user_msg = HumanMessage(
        content=(
            f"{req.message}\n\n"
            f"--- Current document HTML ---\n{current_html[:8000]}"  # cap at 8k chars for context
        )
    )

    response = llm.invoke([system] + history_msgs + [user_msg])
    reply_text = response.content

    # Check if the response contains an updated HTML block
    updated_html: str | None = None
    import re as _re
    html_match = _re.search(r"```html\s*([\s\S]+?)```", reply_text)
    if html_match:
        updated_html = html_match.group(1).strip()
        # Save updated file with draft badge preserved
        updated_html_draft = _inject_draft(updated_html)
        doc_path.write_text(updated_html_draft, encoding="utf-8")
        # Strip code block from reply shown to user
        reply_text = reply_text[:html_match.start()].strip() or "Document updated."

    return ChatResponse(reply=reply_text, updated_html=updated_html)

@app.get("/preview/{filename}", response_class=HTMLResponse)
def preview(filename: str):
    raise HTTPException(status_code=404, detail="No cached preview available")

@app.post("/publish")
def publish_endpoint(payload: dict):
    """Mark a document as published (removes DRAFT badge)."""
    doc_path_rel = payload.get("doc_path", "")
    if not doc_path_rel:
        raise HTTPException(status_code=400, detail="doc_path required")
    doc_path = _DOCS_ROOT / doc_path_rel
    if not doc_path.exists():
        raise HTTPException(status_code=404, detail="Document not found")
    html = doc_path.read_text(encoding="utf-8")
    # Remove the draft badge block injected by _inject_draft
    html = html.replace(DRAFT_BADGE, "")
    doc_path.write_text(html, encoding="utf-8")
    return {"status": "published", "doc_path": doc_path_rel}

@app.post("/api/scrape", response_model=ScrapeResponse)
def scrape_endpoint(req: ScrapeRequest):
    from tools.web_scraper import scrape_and_ingest
    n = scrape_and_ingest(req.url, source=req.source, title=req.title)
    return ScrapeResponse(url=req.url, chunks_ingested=n)


class InsightRequest(BaseModel):
    user_id: str
    insight: str
    namespace: str = "sales_playbook"
    thread_id: str | None = None
    tags: list[str] = []


@app.post("/api/insights")
def create_insight(req: InsightRequest):
    from core.store import record_user_preference

    if not req.user_id or not req.insight:
        raise HTTPException(status_code=400, detail="user_id and insight are required")

    try:
        insight_id = record_user_preference(
            user_id=req.user_id,
            namespace=req.namespace,
            insight=req.insight,
            thread_id=req.thread_id,
            tags=req.tags,
        )
        return {"status": "ok", "insight_id": insight_id}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))


@app.get("/api/insights")
def get_insights(user_id: str, namespace: str = "sales_playbook", query: str = "", top_k: int = 5):
    from core.store import query_user_insights

    if not user_id:
        raise HTTPException(status_code=400, detail="user_id is required")

    try:
        data = query_user_insights(user_id=user_id, namespace=namespace, query=query, top_k=top_k)
        return {"insights": data}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))


@app.get("/api/docs")
def list_docs():
    """Return all HTML documents under CLIENTS/, COMPETITION/, PARTNERS/."""
    _IGNORE = {"index.html", "viewer.html"}
    _SECTION_MAP = {"clients": "client", "competition": "competition", "partners": "partner"}
    docs = []
    docs_root_resolved = _DOCS_ROOT.resolve()
    for section_dir in ["CLIENTS", "COMPETITION", "PARTNERS"]:
        base = _DOCS_ROOT / section_dir
        if not base.exists():
            continue
        for path in sorted(base.rglob("*.html")):
            if path.name in _IGNORE:
                continue
            rel = str(path.relative_to(_DOCS_ROOT))
            parts = rel.replace("\\", "/").split("/")
            section = _SECTION_MAP.get(parts[0].lower(), parts[0].lower())
            category = parts[1] if len(parts) > 2 else ""
            stat = path.stat()
            name = (
                path.stem
                .replace("-presales", "")
                .replace("-", " ")
                .replace("_", " ")
                .title()
                .strip()
            )
            try:
                content = path.read_text(encoding="utf-8", errors="ignore")
                is_draft = DRAFT_BADGE.strip()[:40] in content or "id=\"draft-badge\"" in content
            except Exception:
                is_draft = False
            docs.append({
                "path": rel,
                "filename": path.name,
                "section": section,
                "folder": "/".join(parts[:-1]),
                "category": category,
                "name": name,
                "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%b %d, %Y"),
                "size": stat.st_size,
                "is_draft": is_draft,
            })
    return {"docs": docs}


class DeleteRequest(BaseModel):
    doc_path: str

class DeleteFolderRequest(BaseModel):
    folder_path: str

class MoveRequest(BaseModel):
    doc_path: str
    dest_folder: str  # relative path, e.g. "CLIENTS/RUSTAN_GROUP"

@app.post("/api/delete")
def delete_doc(req: DeleteRequest):
    """Delete a generated HTML document. Only allows deletions within CLIENTS/COMPETITION/PARTNERS."""
    if not req.doc_path:
        raise HTTPException(status_code=400, detail="doc_path required")
    # Resolve and validate — must stay within _DOCS_ROOT, must be .html
    try:
        target = (_DOCS_ROOT / req.doc_path).resolve()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    docs_root_resolved = _DOCS_ROOT.resolve()
    if not str(target).startswith(str(docs_root_resolved)):
        raise HTTPException(status_code=403, detail="Access denied")
    if target.suffix.lower() != ".html":
        raise HTTPException(status_code=400, detail="Only HTML files can be deleted")
    # Only allow deletion inside the three managed sections
    rel_parts = target.relative_to(docs_root_resolved).parts
    if not rel_parts or rel_parts[0].upper() not in ("CLIENTS", "COMPETITION", "PARTNERS"):
        raise HTTPException(status_code=403, detail="Access denied")
    if not target.exists():
        raise HTTPException(status_code=404, detail="File not found")
    target.unlink()
    # Remove parent folder if now empty
    try:
        target.parent.rmdir()
    except OSError:
        pass
    logger.info("[delete] Removed %s", target)
    return {"status": "deleted", "doc_path": req.doc_path}


@app.get("/api/folders")
def list_folders():
    """Return all folder paths (relative to DOCS_ROOT) that currently exist on disk."""
    folders = []
    for section in ["CLIENTS", "COMPETITION", "PARTNERS"]:
        base = _DOCS_ROOT / section
        if not base.exists():
            continue
        for item in sorted(base.rglob("*")):
            if item.is_dir():
                folders.append(str(item.relative_to(_DOCS_ROOT)).replace("\\", "/"))
    return {"folders": folders}


@app.post("/api/delete-folder")
def delete_folder(req: DeleteFolderRequest):
    """Delete a folder and all its HTML contents. Only within CLIENTS/COMPETITION/PARTNERS."""
    if not req.folder_path:
        raise HTTPException(status_code=400, detail="folder_path required")
    try:
        target = (_DOCS_ROOT / req.folder_path).resolve()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    docs_root_resolved = _DOCS_ROOT.resolve()
    if not str(target).startswith(str(docs_root_resolved)):
        raise HTTPException(status_code=403, detail="Access denied")
    rel_parts = target.relative_to(docs_root_resolved).parts
    if not rel_parts or rel_parts[0].upper() not in ("CLIENTS", "COMPETITION", "PARTNERS"):
        raise HTTPException(status_code=403, detail="Access denied")
    if len(rel_parts) < 2:
        raise HTTPException(status_code=403, detail="Cannot delete top-level section")
    if not target.exists() or not target.is_dir():
        raise HTTPException(status_code=404, detail="Folder not found")
    import shutil
    shutil.rmtree(target)
    logger.info("[delete-folder] Removed %s", target)
    return {"status": "deleted", "folder_path": req.folder_path}

@app.post("/api/move")
def move_doc(req: MoveRequest):
    """Move a document to a different folder."""
    if not req.doc_path or not req.dest_folder:
        raise HTTPException(status_code=400, detail="doc_path and dest_folder required")
    try:
        src = (_DOCS_ROOT / req.doc_path).resolve()
        dest_dir = (_DOCS_ROOT / req.dest_folder).resolve()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    docs_root_resolved = _DOCS_ROOT.resolve()
    if not str(src).startswith(str(docs_root_resolved)) or not str(dest_dir).startswith(str(docs_root_resolved)):
        raise HTTPException(status_code=403, detail="Access denied")
    for p in [src, dest_dir]:
        rel = pathlib.Path(p).relative_to(docs_root_resolved).parts
        if not rel or rel[0].upper() not in ("CLIENTS", "COMPETITION", "PARTNERS"):
            raise HTTPException(status_code=403, detail="Access denied")
    if src.suffix.lower() != ".html":
        raise HTTPException(status_code=400, detail="Only HTML files can be moved")
    if not src.exists():
        raise HTTPException(status_code=404, detail="Source not found")
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest_file = dest_dir / src.name
    if dest_file.exists():
        raise HTTPException(status_code=409, detail="A file with the same name already exists in the destination")
    src.rename(dest_file)
    # Clean up empty source folder
    try:
        src.parent.rmdir()
    except OSError:
        pass
    new_path = str(dest_file.relative_to(docs_root_resolved)).replace("\\", "/")
    logger.info("[move] %s -> %s", req.doc_path, new_path)
    return {"status": "moved", "new_path": new_path}

@app.post("/api/crawl")
def crawl_endpoint(req: CrawlRequest):
    from tools.web_scraper import crawl_links
    report = crawl_links(
        req.start_url,
        depth=req.depth,
        url_filter=req.url_filter,
        source=req.source,
        max_pages=req.max_pages,
    )
    return {"pages_scraped": len(report), "details": report}


@app.post("/api/analyse-brand")
def analyse_brand_endpoint(req: BrandAnalyseRequest):
    """
    Fetch the prospect's website and extract:
      - Two dominant brand colours (from CSS custom properties, inline style, or most-frequent hex)
      - Logo URL (og:image, apple-touch-icon, or prominent <img> in header/nav)
    Returns: { colors: ["#hex1", "#hex2"], logo_url: str | null }
    """
    import re as _re
    import httpx
    from urllib.parse import urljoin, urlparse

    url = req.url.strip()
    if not url.startswith(("http://", "https://")):
        url = "https://" + url

    headers = {"User-Agent": "TalonOnePresalesBot/1.0 (+https://talon.one)"}
    try:
        resp = httpx.get(url, headers=headers, timeout=15, follow_redirects=True)
        resp.raise_for_status()
        html = resp.text
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Could not fetch URL: {exc}")

    base = f"{urlparse(url).scheme}://{urlparse(url).netloc}"
    colors: list[str] = []
    logo_url: str | None = None

    # ── 1. Colours: CSS custom properties first ─────────────────────────────
    css_vars = _re.findall(
        r'--(?:color|brand|primary|secondary|accent)[\w-]*\s*:\s*(#[0-9a-fA-F]{6})',
        html, _re.IGNORECASE
    )
    for c in css_vars:
        if c.lower() not in [x.lower() for x in colors]:
            colors.append(c)
        if len(colors) >= 2:
            break

    # ── 2. Colours: background-color / color in <style> blocks ──────────────
    if len(colors) < 2:
        style_blocks = _re.findall(r'<style[^>]*>(.*?)</style>', html, _re.S | _re.I)
        style_text = ' '.join(style_blocks)
        priority_rules = _re.findall(
            r'(?:body|header|nav|\.btn|\.button|\.hero)[^{]*\{[^}]*(?:background(?:-color)?|color)\s*:\s*(#[0-9a-fA-F]{6})',
            style_text, _re.I
        )
        for c in priority_rules:
            if c.lower() in ('#ffffff', '#000000'):
                continue
            if c.lower() not in [x.lower() for x in colors]:
                colors.append(c)
            if len(colors) >= 2:
                break

    # ── 3. Colours: most frequent non-trivial hex in entire HTML ────────────
    if len(colors) < 2:
        skip = {'ffffff', '000000', 'eeeeee', 'f0f0f0', 'cccccc', 'aaaaaa', '999999', 'f4f4f4', 'e0e0e0'}
        freq: dict[str, int] = {}
        for m in _re.finditer(r'#([0-9a-fA-F]{6})(?![0-9a-fA-F])', html):
            h = m.group(1).lower()
            if h not in skip:
                freq[h] = freq.get(h, 0) + 1
        for h, _ in sorted(freq.items(), key=lambda x: -x[1]):
            candidate = '#' + h
            if candidate.lower() not in [x.lower() for x in colors]:
                colors.append(candidate)
            if len(colors) >= 2:
                break

    # ── 4. Logo: og:image ────────────────────────────────────────────────────
    og = _re.search(r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\']([^"\']+)', html, _re.I)
    if not og:
        og = _re.search(r'<meta[^>]+content=["\']([^"\']+)["\'][^>]+property=["\']og:image', html, _re.I)
    if og:
        logo_url = urljoin(base, og.group(1))

    # ── 5. Logo: apple-touch-icon ────────────────────────────────────────────
    if not logo_url:
        touch = _re.search(r'<link[^>]+rel=["\']apple-touch-icon["\'][^>]+href=["\']([^"\']+)', html, _re.I)
        if touch:
            logo_url = urljoin(base, touch.group(1))

    # ── 6. Logo: SVG/PNG in <header>/<nav> ──────────────────────────────────
    if not logo_url:
        nav_block = _re.search(r'<(?:header|nav)[^>]*>(.*?)</(?:header|nav)>', html, _re.S | _re.I)
        if nav_block:
            img = _re.search(
                r'<img[^>]+src=["\']([^"\']+\.(?:svg|png|webp)[^"\']*)["\']',
                nav_block.group(1), _re.I
            )
            if img:
                logo_url = urljoin(base, img.group(1))

    # ── 7. Logo: favicon fallback ────────────────────────────────────────────
    if not logo_url:
        fav = _re.search(r'<link[^>]+rel=["\'][^"\']*(?:icon|shortcut)[^"\']*["\'][^>]+href=["\']([^"\']+)', html, _re.I)
        if fav:
            logo_url = urljoin(base, fav.group(1))

    return {"colors": colors[:2], "logo_url": logo_url}


# ─────────────────────────────────────────────────────────────
# Knowledge Base Management API
# ─────────────────────────────────────────────────────────────

class KBIngestRequest(BaseModel):
    url: str
    source: str          # e.g. "talon_docs", "competitor", "partner", "product_docs"
    label: str = ""      # human-readable label shown in the UI
    topic: str = ""      # knowledge topic category (stored in metadata)
    crawl: bool = False  # if True, follow internal links up to crawl_depth
    crawl_depth: int = 2
    crawl_max_pages: int = 20
    url_filter: str = ""  # optional substring to restrict crawl (e.g. "/api/")

class KBDeleteSourceRequest(BaseModel):
    source: str
    source_id: str = ""  # if set, delete only this one entry; else delete entire source

class KBSearchRequest(BaseModel):
    query: str
    source: str = ""     # filter by source, or empty for all
    top_k: int = 5


@app.get("/api/knowledge/sources")
def kb_list_sources():
    """Return all distinct sources in the knowledge base with chunk counts and last-ingested date."""
    import psycopg
    from psycopg.rows import dict_row
    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")
    try:
        with psycopg.connect(dsn, row_factory=dict_row) as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT
                        source,
                        COUNT(*) AS chunks,
                        COUNT(DISTINCT source_id) AS pages,
                        MAX(created_at) AS last_ingested
                    FROM documents
                    GROUP BY source
                    ORDER BY source
                """)
                rows = cur.fetchall()
        sources = []
        for r in rows:
            sources.append({
                "source": r["source"],
                "chunks": r["chunks"],
                "pages": r["pages"],
                "last_ingested": r["last_ingested"].isoformat() if r["last_ingested"] else None,
            })
        return {"sources": sources}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DB error: {exc}")


@app.get("/api/knowledge/entries")
def kb_list_entries(source: str = "", topic: str = "", limit: int = 100):
    """Return individual document entries (distinct source_ids) for a given source and topic."""
    import psycopg
    from psycopg.rows import dict_row
    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")
    try:
        with psycopg.connect(dsn, row_factory=dict_row) as conn:
            with conn.cursor() as cur:
                if source and topic:
                    cur.execute("""
                        SELECT source, source_id, title, COUNT(*) AS chunks,
                               MIN(created_at) AS created_at, MAX(created_at) AS updated_at,
                               MAX(metadata->>'topic') AS topic
                        FROM documents
                        WHERE source = %s AND metadata->>'topic' = %s
                        GROUP BY source, source_id, title
                        ORDER BY updated_at DESC LIMIT %s
                    """, (source, topic, limit))
                elif source:
                    cur.execute("""
                        SELECT source, source_id, title, COUNT(*) AS chunks,
                               MIN(created_at) AS created_at, MAX(created_at) AS updated_at,
                               MAX(metadata->>'topic') AS topic
                        FROM documents
                        WHERE source = %s
                        GROUP BY source, source_id, title
                        ORDER BY updated_at DESC LIMIT %s
                    """, (source, limit))
                elif topic:
                    cur.execute("""
                        SELECT source, source_id, title, COUNT(*) AS chunks,
                               MIN(created_at) AS created_at, MAX(created_at) AS updated_at,
                               MAX(metadata->>'topic') AS topic
                        FROM documents
                        WHERE metadata->>'topic' = %s
                        GROUP BY source, source_id, title
                        ORDER BY updated_at DESC LIMIT %s
                    """, (topic, limit))
                else:
                    cur.execute("""
                        SELECT source, source_id, title, COUNT(*) AS chunks,
                               MIN(created_at) AS created_at, MAX(created_at) AS updated_at,
                               MAX(metadata->>'topic') AS topic
                        FROM documents
                        GROUP BY source, source_id, title
                        ORDER BY updated_at DESC LIMIT %s
                    """, (limit,))
                rows = cur.fetchall()
        entries = [{
            "source": r["source"],
            "source_id": r["source_id"],
            "title": r["title"],
            "chunks": r["chunks"],
            "topic": r["topic"],
            "created_at": r["created_at"].isoformat() if r["created_at"] else None,
            "updated_at": r["updated_at"].isoformat() if r["updated_at"] else None,
        } for r in rows]
        return {"entries": entries}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DB error: {exc}")


@app.post("/api/knowledge/ingest")
def kb_ingest(req: KBIngestRequest):
    """
    Ingest a URL into the knowledge base with a given source label.
    If crawl=True, follows internal links up to crawl_depth.
    """
    from tools.web_scraper import scrape_and_ingest, crawl_links
    try:
        if req.crawl:
            report = crawl_links(
                req.start_url if hasattr(req, "start_url") else req.url,
                depth=req.crawl_depth,
                url_filter=req.url_filter,
                source=req.source,
                max_pages=req.crawl_max_pages,
            )
            total = sum(r.get("chunks_ingested", 0) for r in report if isinstance(r, dict))
            return {
                "status": "ok",
                "url": req.url,
                "source": req.source,
                "pages_crawled": len(report),
                "chunks_ingested": total,
            }
        else:
            n = scrape_and_ingest(req.url, source=req.source, title=req.label or req.url)
            return {
                "status": "ok",
                "url": req.url,
                "source": req.source,
                "chunks_ingested": n,
            }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))


@app.delete("/api/knowledge/source")
def kb_delete_source(req: KBDeleteSourceRequest):
    """Delete all chunks for a source, or a specific source_id entry."""
    import psycopg
    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")
    if not req.source:
        raise HTTPException(status_code=400, detail="source is required")
    try:
        with psycopg.connect(dsn) as conn:
            with conn.cursor() as cur:
                if req.source_id:
                    cur.execute(
                        "DELETE FROM documents WHERE source = %s AND source_id LIKE %s",
                        (req.source, f"{req.source_id}%"),
                    )
                else:
                    cur.execute("DELETE FROM documents WHERE source = %s", (req.source,))
                deleted = cur.rowcount
                conn.commit()
        return {"status": "deleted", "source": req.source, "chunks_deleted": deleted}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DB error: {exc}")


@app.post("/api/knowledge/search")
def kb_search(req: KBSearchRequest):
    """Test a semantic search query against the knowledge base."""
    from tools.postgres_rag import similarity_search
    import psycopg
    from psycopg.rows import dict_row
    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")
    try:
        if req.source:
            # Source-filtered search via raw SQL
            from tools.postgres_rag import _embed
            vector = _embed(req.query)
            vector_literal = "[" + ",".join(str(v) for v in vector) + "]"
            with psycopg.connect(dsn, row_factory=dict_row) as conn:
                with conn.cursor() as cur:
                    cur.execute(
                        f"""
                        SELECT source, source_id, title, content,
                               1 - (embedding <=> %s::vector) AS score
                        FROM documents
                        WHERE source = %s
                        ORDER BY embedding <=> %s::vector
                        LIMIT %s
                        """,
                        (vector_literal, req.source, vector_literal, req.top_k),
                    )
                    rows = cur.fetchall()
            results = [{
                "source": r["source"], "source_id": r["source_id"],
                "title": r["title"], "score": round(float(r["score"]), 4),
                "excerpt": r["content"][:300],
            } for r in rows]
        else:
            chunks = similarity_search(req.query, top_k=req.top_k)
            results = [{
                "source": c["source"], "source_id": c["source_id"],
                "score": round(float(c["score"]), 4),
                "excerpt": c["content"][:300],
            } for c in chunks]
        return {"query": req.query, "results": results}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))


# ─────────────────────────────────────────────────────────────
# Knowledge Base — Additional Ingestion Endpoints
# ─────────────────────────────────────────────────────────────

class KBGDriveRequest(BaseModel):
    url: str
    source: str = "google_drive"
    title: str = ""
    topic: str = ""
    file_type: str = "auto"   # auto | document | spreadsheets | presentation | pdf

class KBNotionRequest(BaseModel):
    url: str
    source: str = "notion"
    title: str = ""
    topic: str = ""
    notion_token: str = ""    # override NOTION_API_KEY for per-request auth

class KBAnalyseRequest(BaseModel):
    source: str
    source_id: str = ""
    content_sample: str = ""  # optional: pass content directly instead of fetching from DB
    refine_prompt: str = ""

class KBSaveInsightRequest(BaseModel):
    source: str
    source_id: str
    topic: str = ""
    subjects: list[str] = []
    key_insights: list[dict] = []
    value_propositions: list[str] = []
    competitive_advantages: list[str] = []
    diagram_suggestions: list[dict] = []
    presales_notes: str = ""


def _ingest_text_chunks(
    text: str,
    source: str,
    source_id: str,
    title: str = "",
    topic: str = "",
    chunk_size: int = 800,
    overlap: int = 80,
) -> int:
    """Split text and ingest chunks into pgvector. Returns chunk count."""
    from tools.postgres_rag import ingest_document, chunk_text
    chunks = chunk_text(text, size=chunk_size, overlap=overlap)
    metadata = {"topic": topic} if topic else {}
    for chunk in chunks:
        ingest_document(chunk, source=source, source_id=source_id, title=title, metadata=metadata)
    return len(chunks)


@app.post("/api/knowledge/upload")
async def kb_upload(
    file: UploadFile = File(...),
    source: str = Form("web"),
    title: str = Form(""),
    topic: str = Form(""),
):
    """Upload a document file (PDF, DOCX, PPTX, XLSX, TXT, CSV) and ingest into KB."""
    import io

    content_bytes = await file.read()
    MAX_SIZE = 50 * 1024 * 1024
    if len(content_bytes) > MAX_SIZE:
        raise HTTPException(status_code=413, detail="File too large — max 50 MB")

    filename = file.filename or "upload"
    ext = os.path.splitext(filename)[1].lower()
    file_title = title or filename
    text = ""

    if ext == ".pdf":
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(content_bytes))
            pages = [p.extract_text() or "" for p in reader.pages]
            text = "\n\n".join(p for p in pages if p.strip())
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"PDF extraction failed: {exc}")

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if parts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides)
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"PPTX extraction failed: {exc}")

    elif ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            sheets = []
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
            text = "\n\n".join(sheets)
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"XLSX extraction failed: {exc}")

    elif ext == ".docx":
        try:
            try:
                import docx as _docx
                doc = _docx.Document(io.BytesIO(content_bytes))
                text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                # Fallback: extract from raw XML
                import zipfile, xml.etree.ElementTree as ET
                with zipfile.ZipFile(io.BytesIO(content_bytes)) as z:
                    with z.open("word/document.xml") as f:
                        tree = ET.parse(f)
                ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                paragraphs = ["".join(t.text or "" for t in para.findall(".//w:t", ns))
                              for para in tree.findall(".//w:p", ns)]
                text = "\n\n".join(p for p in paragraphs if p.strip())
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"DOCX extraction failed: {exc}")

    elif ext in (".txt", ".csv", ".md"):
        text = content_bytes.decode("utf-8", errors="ignore")

    else:
        raise HTTPException(status_code=422, detail=f"Unsupported file type: {ext}")

    if not text.strip():
        raise HTTPException(status_code=422, detail="No text could be extracted from the file")

    source_id = f"upload:{filename}"
    n = _ingest_text_chunks(text, source=source, source_id=source_id, title=file_title, topic=topic)
    logger.info("[upload] Ingested %s → %d chunks (source=%s)", filename, n, source)
    return {"status": "ok", "filename": filename, "source": source, "source_id": source_id, "chunks_ingested": n}


@app.post("/api/knowledge/gdrive")
def kb_gdrive(req: KBGDriveRequest):
    """
    Ingest from a Google Drive URL.
    Supports: Google Docs, Sheets, Slides (via export URL) and binary files (PDF/PPTX).
    The file/page must be accessible to anyone with the link.
    """
    import re as _re
    import io
    import httpx
    from urllib.parse import urlparse

    url = req.url.strip()
    headers = {"User-Agent": "TalonOnePresalesBot/1.0"}

    # ── Detect file type and extract file/doc ID ──────────────────────────
    file_type = req.file_type
    file_id = None

    # Google Docs: docs.google.com/document/d/FILE_ID/...
    m = _re.search(r"docs\.google\.com/(document|spreadsheets|presentation)/d/([^/\?]+)", url)
    if m:
        doc_type = m.group(1)  # document | spreadsheets | presentation
        file_id = m.group(2)
        if file_type == "auto":
            file_type = doc_type

    # Drive file: drive.google.com/file/d/FILE_ID/...
    if not file_id:
        m = _re.search(r"drive\.google\.com/file/d/([^/\?]+)", url)
        if m:
            file_id = m.group(1)
            if file_type == "auto":
                file_type = "pdf"

    # drive.google.com/open?id=FILE_ID or /uc?id=FILE_ID
    if not file_id:
        m = _re.search(r"[?&]id=([a-zA-Z0-9_-]+)", url)
        if m:
            file_id = m.group(1)

    if not file_id:
        raise HTTPException(status_code=400, detail="Could not extract a file ID from the provided URL. Make sure it is a valid Google Drive or Google Docs link.")

    # ── Build export/download URL ─────────────────────────────────────────
    if file_type in ("document",):
        export_url = f"https://docs.google.com/document/d/{file_id}/export?format=txt"
    elif file_type in ("spreadsheets",):
        export_url = f"https://docs.google.com/spreadsheets/d/{file_id}/export?format=csv"
    elif file_type in ("presentation",):
        export_url = f"https://docs.google.com/{file_type}/d/{file_id}/export?format=txt"
    else:
        # Binary download
        export_url = f"https://drive.google.com/uc?export=download&id={file_id}"

    try:
        resp = httpx.get(export_url, headers=headers, timeout=30, follow_redirects=True)
        # GDrive sometimes returns a confirmation page for large files
        if resp.status_code == 200 and b"virus scan warning" in resp.content[:500].lower():
            # Extract the confirm token and retry
            token_m = _re.search(rb'confirm=([a-zA-Z0-9_-]+)', resp.content)
            if token_m:
                token = token_m.group(1).decode()
                export_url += f"&confirm={token}"
                resp = httpx.get(export_url, headers=headers, timeout=30, follow_redirects=True)
        resp.raise_for_status()
    except httpx.HTTPStatusError as exc:
        raise HTTPException(status_code=502, detail=f"Could not fetch from Google Drive ({exc.response.status_code}). Make sure the file is shared as 'Anyone with the link'.")
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Could not fetch from Google Drive: {exc}")

    # ── Extract text ──────────────────────────────────────────────────────
    text = ""
    content_type = resp.headers.get("content-type", "")

    if file_type == "pdf" or "pdf" in content_type:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(resp.content))
            text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"PDF extraction failed: {exc}")
    elif file_type == "presentation":
        # Exported as .pptx binary
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(resp.content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if parts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides)
        except Exception:
            text = resp.text  # fallback to raw text
    else:
        text = resp.text

    if not text.strip():
        raise HTTPException(status_code=422, detail="No text could be extracted from the Drive file")

    source_id = f"gdrive:{file_id}"
    file_title = req.title or source_id
    n = _ingest_text_chunks(text, source=req.source, source_id=source_id, title=file_title, topic=req.topic)
    logger.info("[gdrive] Ingested %s → %d chunks (source=%s)", source_id, n, req.source)
    return {"status": "ok", "source_id": source_id, "source": req.source, "chunks_ingested": n}


@app.get("/api/knowledge/notion-pages")
def kb_notion_pages(q: str = "", limit: int = 30):
    """
    Search the connected Notion workspace and return matching pages.
    Call with q='' to browse all recent pages.
    """
    try:
        from tools.notion import search_notion
        pages = search_notion(q, max_pages=limit)
        return {"pages": pages, "total": len(pages)}
    except EnvironmentError as exc:
        raise HTTPException(status_code=503, detail=str(exc))
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Notion search failed: {exc}")


@app.post("/api/knowledge/notion-ingest")
def kb_notion_ingest(req: KBNotionRequest):
    """
    Ingest a Notion page by URL. Uses NOTION_API_KEY from .env by default;
    pass notion_token to override for a specific request.
    """
    import re as _re

    url = req.url.strip()

    # Extract page ID from URL
    # Formats: /page-title-PAGE_ID  or  /PAGE_ID  or  notion.so/workspace/title-PAGE_ID
    page_id_match = _re.search(r"([a-f0-9]{32}|[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12})$",
                                url.rstrip("/").split("?")[0])
    if not page_id_match:
        # Try last path segment (may have title-id pattern)
        slug = url.rstrip("/").split("/")[-1].split("?")[0]
        id_part = slug.split("-")[-1]
        if len(id_part) == 32:
            page_id = id_part
        else:
            raise HTTPException(status_code=400, detail="Could not extract a Notion page ID from the URL. Use a direct page URL like https://notion.so/your-workspace/Page-Title-abc123...")
    else:
        page_id = page_id_match.group(1).replace("-", "")

    # Override token if provided
    if req.notion_token:
        os.environ["NOTION_API_KEY"] = req.notion_token

    try:
        from tools.notion import extract_page_text, ingest_notion_page
        # The existing ingest_notion_page function embeds and stores in pgvector
        ingested = ingest_notion_page(page_id)
        source_id = f"notion:{page_id}"
        logger.info("[notion] Ingested page %s → %d chunks", page_id, ingested)
        return {"status": "ok", "source_id": source_id, "source": req.source, "chunks_ingested": ingested}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Notion ingestion failed: {exc}")


@app.post("/api/knowledge/analyse")
def kb_analyse(req: KBAnalyseRequest):
    """
    AI analysis of ingested content.
    Extracts topic, subjects, key insights, value propositions,
    competitive advantages, and diagram suggestions using Gemini.
    Results can be saved via /api/knowledge/save-insight.
    """
    import psycopg
    from psycopg.rows import dict_row
    import json as _json
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_core.messages import SystemMessage, HumanMessage
    from core.prompts import KB_ANALYSE_SYSTEM

    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")

    content_sample = req.content_sample
    if not content_sample and req.source:
        try:
            with psycopg.connect(dsn, row_factory=dict_row) as conn:
                with conn.cursor() as cur:
                    if req.source_id:
                        cur.execute(
                            "SELECT content, title FROM documents WHERE source = %s AND source_id LIKE %s ORDER BY id LIMIT 12",
                            (req.source, req.source_id + "%"),
                        )
                    else:
                        cur.execute(
                            "SELECT content, title FROM documents WHERE source = %s ORDER BY created_at DESC LIMIT 12",
                            (req.source,),
                        )
                    rows = cur.fetchall()
            parts = []
            for r in rows:
                prefix = f"[{r['title']}]\n" if r["title"] else ""
                parts.append(prefix + r["content"])
            content_sample = "\n\n---\n\n".join(parts)[:6000]
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"DB error fetching content: {exc}")

    if not content_sample:
        raise HTTPException(status_code=400, detail="No content found — ingest the source first")

    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        google_api_key=os.getenv("GOOGLE_API_KEY", ""),
    )

    user_msg = f"Please analyse this content:\n\n{content_sample}"
    if req.refine_prompt:
        user_msg += f"\n\nAdditional instructions from the presales expert: {req.refine_prompt}"

    try:
        response = llm.invoke([SystemMessage(content=KB_ANALYSE_SYSTEM), HumanMessage(content=user_msg)])
        raw = response.content.strip()
        # Strip markdown fences if present
        import re as _re
        raw = _re.sub(r"^```(?:json)?\s*", "", raw)
        raw = _re.sub(r"\s*```$", "", raw)
        result = _json.loads(raw)
    except _json.JSONDecodeError:
        # Try to extract JSON from the response
        import re as _re
        json_match = _re.search(r"\{[\s\S]+\}", response.content)
        result = _json.loads(json_match.group()) if json_match else {
            "error": "Could not parse AI response",
            "raw": response.content[:400],
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"AI analysis failed: {exc}")

    return result


def _ensure_insights_table(conn):
    """Create kb_insights table if it doesn't exist."""
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS kb_insights (
                id SERIAL PRIMARY KEY,
                source TEXT NOT NULL,
                source_id TEXT NOT NULL,
                topic TEXT,
                subjects TEXT[],
                key_insights JSONB DEFAULT '[]',
                value_propositions TEXT[],
                competitive_advantages TEXT[],
                diagram_suggestions JSONB DEFAULT '[]',
                presales_notes TEXT,
                created_at TIMESTAMPTZ DEFAULT NOW(),
                updated_at TIMESTAMPTZ DEFAULT NOW()
            )
        """)
        conn.commit()


@app.post("/api/knowledge/save-insight")
def kb_save_insight(req: KBSaveInsightRequest):
    """Save a presales-curated insight derived from KB analysis. Agents use these when generating documents."""
    import psycopg
    import json as _json

    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")

    try:
        with psycopg.connect(dsn) as conn:
            _ensure_insights_table(conn)
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO kb_insights
                        (source, source_id, topic, subjects, key_insights,
                         value_propositions, competitive_advantages, diagram_suggestions, presales_notes)
                    VALUES (%s,%s,%s,%s,%s::jsonb,%s,%s,%s::jsonb,%s)
                    ON CONFLICT DO NOTHING
                """, (
                    req.source, req.source_id, req.topic or None,
                    req.subjects or None,
                    _json.dumps(req.key_insights),
                    req.value_propositions or None,
                    req.competitive_advantages or None,
                    _json.dumps(req.diagram_suggestions),
                    req.presales_notes or None,
                ))
                conn.commit()
        logger.info("[save-insight] Saved insight for %s / %s", req.source, req.source_id)
        return {"status": "saved"}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DB error: {exc}")


@app.get("/api/knowledge/insights")
def kb_list_insights(source: str = "", topic: str = "", limit: int = 50):
    """Return saved presales insights, optionally filtered by source or topic."""
    import psycopg
    from psycopg.rows import dict_row

    dsn = os.getenv("DATABASE_URL", "postgresql+psycopg://presales:presales_secret@localhost:5432/presales_db")
    dsn = dsn.replace("postgresql+psycopg://", "postgresql://")

    try:
        with psycopg.connect(dsn, row_factory=dict_row) as conn:
            _ensure_insights_table(conn)
            with conn.cursor() as cur:
                conditions = []
                params: list = []
                if source:
                    conditions.append("source = %s")
                    params.append(source)
                if topic:
                    conditions.append("topic = %s")
                    params.append(topic)
                where = ("WHERE " + " AND ".join(conditions)) if conditions else ""
                params.append(limit)
                cur.execute(f"""
                    SELECT id, source, source_id, topic, subjects, key_insights,
                           value_propositions, competitive_advantages, diagram_suggestions,
                           presales_notes, created_at
                    FROM kb_insights {where}
                    ORDER BY created_at DESC LIMIT %s
                """, params)
                rows = cur.fetchall()
        return {"insights": [dict(r) for r in rows]}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DB error: {exc}")


def _start_server():
    """Start the uvicorn server."""
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)


# ─────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Talon.One Presales Multi-Agent System")
    parser.add_argument(
        "--prospect",
        metavar="FILE",
        help="Path to a JSON file with prospect fields (skips interactive prompt)",
    )
    parser.add_argument(
        "--serve",
        action="store_true",
        help="Start the FastAPI server instead of running a single presales job",
    )
    parser.add_argument(
        "--resume",
        action="store_true",
        help="Resume from the last checkpoint for this prospect (same company name)",
    )
    parser.add_argument(
        "--history",
        metavar="COMPANY",
        help="Print checkpoint history for a given company name and exit",
    )
    args = parser.parse_args()

    if args.serve:
        _start_server()
        return

    if args.history:
        graph = _import_graph()
        thread_id = make_thread_id(args.history)
        history = get_thread_history(graph, thread_id)
        print(f"\nCheckpoint history for '{args.history}' (thread={thread_id}):")
        for h in history:
            deployed = "DEPLOYED" if h["deployed"] else ("HTML" if h["html_ready"] else "in-progress")
            print(f"  [{h['ts'][:19]}] step={h['step'][:8]}  agent={h['agent']:<20} iter={h['iteration']}  {deployed}")
        if not history:
            print("  No checkpoints found.")
        return

    if args.prospect:
        with open(args.prospect, "r", encoding="utf-8") as fh:
            data = json.load(fh)
        prospect = ProspectContext(**data)
    else:
        prospect = _prompt_prospect()

    final_state = run_agent(prospect, resume=args.resume)

    # Print a brief summary
    print("\n── Result ───────────────────────────────────────────────")
    print(f"Company   : {final_state['prospect'].get('company_name')}")
    print(f"File      : {final_state.get('output_filename', 'N/A')}")
    print(f"Vercel    : {final_state.get('vercel_url', 'N/A')}")
    print(f"GitHub    : {final_state.get('github_pr_url', 'N/A')}")
    v = final_state.get("validation")
    if v:
        status = "✅ PASSED" if v["passed"] else "❌ FAILED"
        print(f"Validation: {status}")
        for note in v["notes"]:
            print(f"  {note}")
    if final_state.get("error"):
        print(f"Error     : {final_state['error']}")
        sys.exit(1)


if __name__ == "__main__":
    main()
