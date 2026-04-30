"""
tools/google_drive.py
──────────────────────
Google Drive integration for the Data Architect agent.

Supports:
  - Listing files inside a folder (with optional keyword filter)
  - Downloading & extracting text from: PDF, PPTX, XLSX, PNG/JPEG (via Vision LLM)

Authentication: OAuth 2.0 via credentials JSON (set GOOGLE_CREDENTIALS_PATH in .env).
The first run opens a browser for the consent flow and caches a token in
./credentials/token.json.
"""

from __future__ import annotations

import io
import logging
import os
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
_TOKEN_PATH = Path("./credentials/token.json")
_CREDENTIALS_PATH = Path(os.getenv("GOOGLE_CREDENTIALS_PATH", "./credentials/google_credentials.json"))


# ─────────────────────────────────────────────────────────────
# Auth
# ─────────────────────────────────────────────────────────────

def _get_drive_service() -> Any:
    """Return an authenticated Google Drive service client."""
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build

    creds = None
    if _TOKEN_PATH.exists():
        creds = Credentials.from_authorized_user_file(str(_TOKEN_PATH), _SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(str(_CREDENTIALS_PATH), _SCOPES)
            creds = flow.run_local_server(port=0)
        _TOKEN_PATH.parent.mkdir(parents=True, exist_ok=True)
        _TOKEN_PATH.write_text(creds.to_json())

    return build("drive", "v3", credentials=creds)


# ─────────────────────────────────────────────────────────────
# File listing
# ─────────────────────────────────────────────────────────────

def list_drive_files(folder_id: str, query: str = "") -> list[dict[str, str]]:
    """
    List files inside a Google Drive folder.

    Args:
        folder_id: The Drive folder ID to search in.
        query: Optional keyword to filter by file name.

    Returns:
        List of dicts with keys: id, name, mimeType.
    """
    service = _get_drive_service()

    q_parts = [f"'{folder_id}' in parents", "trashed = false"]
    if query:
        q_parts.append(f"name contains '{query}'")

    results = (
        service.files()
        .list(
            q=" and ".join(q_parts),
            fields="files(id, name, mimeType)",
            pageSize=20,
        )
        .execute()
    )

    files = results.get("files", [])
    logger.info("[GoogleDrive] Found %d files in folder %s", len(files), folder_id)
    return files


# ─────────────────────────────────────────────────────────────
# Download and extract text
# ─────────────────────────────────────────────────────────────

def download_and_extract_text(file_id: str, mime_type: str) -> str:
    """
    Download a Drive file and return its text content.

    Supported MIME types:
      - application/pdf
      - application/vnd.openxmlformats-officedocument.presentationml.presentation (PPTX)
      - application/vnd.openxmlformats-officedocument.spreadsheetml.sheet (XLSX)
      - image/png, image/jpeg → Vision LLM description

    Google Workspace types (Docs/Slides/Sheets) are exported to their Office equivalents
    before parsing.
    """
    service = _get_drive_service()

    # Google Workspace export targets
    EXPORT_MAP = {
        "application/vnd.google-apps.document": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.google-apps.presentation": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.google-apps.spreadsheet": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }

    if mime_type in EXPORT_MAP:
        export_mime = EXPORT_MAP[mime_type]
        data = service.files().export_media(fileId=file_id, mimeType=export_mime).execute()
        # Update effective mime type for parsing
        mime_type = export_mime
    else:
        data = service.files().get_media(fileId=file_id).execute()

    buf = io.BytesIO(data)

    if mime_type == "application/pdf":
        return _extract_pdf(buf)
    if "presentationml" in mime_type:
        return _extract_pptx(buf)
    if "spreadsheetml" in mime_type:
        return _extract_xlsx(buf)
    if mime_type in ("image/png", "image/jpeg", "image/webp"):
        return _describe_image_via_vision(data)

    # Fallback: try UTF-8 decode
    try:
        return data.decode("utf-8", errors="replace")
    except Exception:
        return ""


# ─────────────────────────────────────────────────────────────
# Format-specific parsers
# ─────────────────────────────────────────────────────────────

def _extract_pdf(buf: io.BytesIO) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(buf)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_pptx(buf: io.BytesIO) -> str:
    from pptx import Presentation

    prs = Presentation(buf)
    texts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts = [f"## Slide {slide_num}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        texts.append("\n".join(parts))
    return "\n\n".join(texts)


def _extract_xlsx(buf: io.BytesIO) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(buf, read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        rows.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            rows.append("\t".join(str(c) if c is not None else "" for c in row))
    return "\n".join(rows)


def _describe_image_via_vision(image_bytes: bytes) -> str:
    """Send an image to the Vision LLM and return a structured description."""
    import base64
    import os

    model = os.getenv("LLM_MODEL", "gemini-2.5-pro")

    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_core.messages import HumanMessage

    b64 = base64.b64encode(image_bytes).decode()
    llm = ChatGoogleGenerativeAI(model=model)
    response = llm.invoke(
        [HumanMessage(content=[
            {"type": "text", "text": "Describe this architecture diagram in detail, focusing on integration flows, components, and data direction."},
            {"type": "image_url", "image_url": f"data:image/png;base64,{b64}"},
        ])]
    )

    return response.content if hasattr(response, "content") else str(response)
